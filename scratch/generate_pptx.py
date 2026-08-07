import sys
import os

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE
except ImportError:
    print("python-pptx is not installed. Please run: pip install python-pptx")
    sys.exit(1)

def create_presentation():
    prs = Presentation()
    
    # Set slide dimensions to widescreen (16:9)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Color palette
    DARK_BG = RGBColor(15, 23, 42)      # #0f172a Deep Slate
    CARD_BG = RGBColor(30, 41, 59)      # #1e293b Charcoal
    KKU_ORANGE = RGBColor(249, 115, 22) # #f97316 KKU Orange
    WHITE = RGBColor(248, 250, 252)     # #f8fafc Primary Text
    MUTED_GRAY = RGBColor(148, 163, 184)# #94a3b8 Secondary Text
    GREEN = RGBColor(16, 185, 129)      # #10b981 Success Green

    def apply_dark_background(slide):
        # Create a full-slide rectangle for background color
        background = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
        )
        background.fill.solid()
        background.fill.fore_color.rgb = DARK_BG
        background.line.color.rgb = DARK_BG # borderless
        # Send to back by moving it to the first element in shape tree
        slide.shapes._spTree.remove(background._element)
        slide.shapes._spTree.insert(2, background._element)

    def add_slide_header(slide, title, subtitle):
        apply_dark_background(slide)
        
        # Subtitle (Top Small)
        sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.4))
        tf_sub = sub_box.text_frame
        tf_sub.word_wrap = True
        tf_sub.margin_left = tf_sub.margin_top = tf_sub.margin_bottom = tf_sub.margin_right = 0
        p_sub = tf_sub.paragraphs[0]
        p_sub.text = subtitle.upper()
        p_sub.font.name = 'Sarabun'
        p_sub.font.size = Pt(12)
        p_sub.font.bold = True
        p_sub.font.color.rgb = KKU_ORANGE
        
        # Main Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.8))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        tf_title.margin_left = tf_title.margin_top = tf_title.margin_bottom = tf_title.margin_right = 0
        p_title = tf_title.paragraphs[0]
        p_title.text = title
        p_title.font.name = 'Sarabun'
        p_title.font.size = Pt(32)
        p_title.font.bold = True
        p_title.font.color.rgb = WHITE
        
        # Divider Line
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.6), Inches(11.733), Inches(0.02))
        line.fill.solid()
        line.fill.fore_color.rgb = CARD_BG
        line.line.fill.background()

    # ==================== SLIDE 1: COVER ====================
    slide_layout = prs.slide_layouts[6] # Blank
    slide1 = prs.slides.add_slide(slide_layout)
    apply_dark_background(slide1)
    
    # Decorative Top border
    top_bar = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.1))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = KKU_ORANGE
    top_bar.line.fill.background()

    # Logo Placeholder shape (Hexagon)
    logo = slide1.shapes.add_shape(MSO_SHAPE.HEXAGON, Inches(5.9), Inches(1.2), Inches(1.5), Inches(1.5))
    logo.fill.solid()
    logo.fill.fore_color.rgb = CARD_BG
    logo.line.color.rgb = KKU_ORANGE
    logo.line.width = Pt(3)

    # Main Title Box
    title_box = slide1.shapes.add_textbox(Inches(1.0), Inches(3.2), Inches(11.333), Inches(2.2))
    tf1 = title_box.text_frame
    tf1.word_wrap = True
    
    p1 = tf1.paragraphs[0]
    p1.text = "ระบบแดชบอร์ดประเมินและติดตามความปลอดภัยทางไซเบอร์ มข."
    p1.alignment = PP_ALIGN.CENTER
    p1.font.name = 'Sarabun'
    p1.font.size = Pt(34)
    p1.font.bold = True
    p1.font.color.rgb = WHITE
    
    p2 = tf1.add_paragraph()
    p2.text = "KKU Security Scorecard & Vulnerability Management Dashboard"
    p2.alignment = PP_ALIGN.CENTER
    p2.font.name = 'Sarabun'
    p2.font.size = Pt(18)
    p2.font.color.rgb = MUTED_GRAY
    p2.space_before = Pt(10)

    # Meta Info
    meta_box = slide1.shapes.add_textbox(Inches(1.0), Inches(5.8), Inches(11.333), Inches(1.0))
    tf_meta = meta_box.text_frame
    tf_meta.word_wrap = True
    p_meta = tf_meta.paragraphs[0]
    p_meta.text = "คณะผู้พัฒนาโครงงาน มหาวิทยาลัยขอนแก่น  |  โครงงานเทคโนโลยีสารสนเทศ ปีการศึกษา 2569"
    p_meta.alignment = PP_ALIGN.CENTER
    p_meta.font.name = 'Sarabun'
    p_meta.font.size = Pt(14)
    p_meta.font.color.rgb = KKU_ORANGE

    # ==================== SLIDE 2: WHAT IS IT ====================
    slide2 = prs.slides.add_slide(slide_layout)
    add_slide_header(slide2, "ทำอะไร (What is it?)", "01. Overview")
    
    # Highlight Box
    hb = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2.0), Inches(11.733), Inches(1.0))
    hb.fill.solid()
    hb.fill.fore_color.rgb = CARD_BG
    hb.line.color.rgb = KKU_ORANGE
    hb.line.width = Pt(1.5)
    tf_hb = hb.text_frame
    tf_hb.word_wrap = True
    tf_hb.margin_left = Inches(0.3)
    p_hb = tf_hb.paragraphs[0]
    p_hb.text = "KKU Security Scorecard คือ ระบบแดชบอร์ดประเมินและวิเคราะห์ความปลอดภัยไซเบอร์แบบ Full-stack ที่รวบรวมข้อมูลสถานะช่องโหว่ความเสี่ยงของทุกคณะและหน่วยงานในมหาวิทยาลัยขอนแก่นมาแสดงผลในที่เดียว"
    p_hb.font.name = 'Sarabun'
    p_hb.font.size = Pt(16)
    p_hb.font.color.rgb = WHITE

    # 3 Grid Cards
    card_width = Inches(3.64)
    card_height = Inches(3.3)
    card_gap = Inches(0.4)
    y_pos = Inches(3.4)
    
    cards_data = [
        ("Security Rating", "ประเมินและวัดผลระบบเครือข่าย โดเมน และเว็บไซต์ โดยตัดเกรดเป็นตัวอักษร (A-F) เพื่อให้เห็นระดับความปลอดภัยของแต่ละหน่วยงานอย่างชัดเจนและเข้าใจง่าย"),
        ("Centralized DB", "รวบรวมและจัดเก็บข้อมูลช่องโหว่และคะแนนความเสี่ยงจาก API และการอัปโหลดไฟล์ (CSV) ช่วยในการประมวลผล เปรียบเทียบข้อมูลย้อนหลัง และวิเคราะห์แนวโน้ม"),
        ("Interactive Map", "แสดงพิกัดที่ตั้งคณะย่อยผ่าน Interactive Campus Map โดยปักหมุดสีตามเกรดความปลอดภัย เพื่อช่วยในการชี้เป้าหน่วยงานที่มีความเสี่ยงสูงได้ทันที")
    ]
    
    for i, (title, desc) in enumerate(cards_data):
        x_pos = Inches(0.8) + i * (card_width + card_gap)
        card = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x_pos, y_pos, card_width, card_height)
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = DARK_BG
        
        tf_card = card.text_frame
        tf_card.word_wrap = True
        tf_card.margin_left = Inches(0.2)
        tf_card.margin_right = Inches(0.2)
        tf_card.margin_top = Inches(0.2)
        
        p_c1 = tf_card.paragraphs[0]
        p_c1.text = title
        p_c1.font.name = 'Sarabun'
        p_c1.font.size = Pt(18)
        p_c1.font.bold = True
        p_c1.font.color.rgb = KKU_ORANGE
        p_c1.space_after = Pt(12)
        
        p_c2 = tf_card.add_paragraph()
        p_c2.text = desc
        p_c2.font.name = 'Sarabun'
        p_c2.font.size = Pt(13)
        p_c2.font.color.rgb = MUTED_GRAY

    # ==================== SLIDE 3: WHY ====================
    slide3 = prs.slides.add_slide(slide_layout)
    add_slide_header(slide3, "ทำไมถึงทำ (Why did we do it?)", "02. Motivation")
    
    # Left Content (Bullets)
    left_box = slide3.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(6.5), Inches(4.5))
    tf_l = left_box.text_frame
    tf_l.word_wrap = True
    
    bullets = [
        ("โครงข่ายขนาดใหญ่และซับซ้อน: ", "มหาวิทยาลัยขอนแก่นมีคณะย่อยและหน่วยงานกว่า 21 แห่ง ซึ่งดูแลรับผิดชอบเซิร์ฟเวอร์และอุปกรณ์เครือข่ายของตนเอง แยกจากส่วนกลาง"),
        ("ภัยคุกคามรอบด้าน: ", "ความเสี่ยงจากเซิร์ฟเวอร์ที่ละเลยการติดตั้ง Patch (CVE), โดเมนหลอกลวง, DNS Spoofing, และความปลอดภัยของเว็บไซต์ย่อย"),
        ("ขาดศูนย์กลางข้อมูล (No Visibility): ", "ผู้บริหารและทีมความปลอดภัยส่วนกลางขาดหน้าจอ Dashboard บูรณาการคะแนนสุขภาพของทุกคณะพร้อมกัน"),
        ("การทำงานและแก้ไขที่ล่าช้า: ", "เมื่อระบบตรวจพบช่องโหว่ การชี้เป้าเครื่องเซิร์ฟเวอร์และระบุตัวผู้ดูแลระบบของแต่ละหน่วยงานย่อยใช้เวลาทำงานนาน")
    ]
    
    for i, (bold_txt, regular_txt) in enumerate(bullets):
        p = tf_l.add_paragraph() if i > 0 else tf_l.paragraphs[0]
        p.space_after = Pt(16)
        
        run1 = p.add_run()
        run1.text = "•  " + bold_txt
        run1.font.name = 'Sarabun'
        run1.font.size = Pt(14)
        run1.font.bold = True
        run1.font.color.rgb = KKU_ORANGE
        
        run2 = p.add_run()
        run2.text = regular_txt
        run2.font.name = 'Sarabun'
        run2.font.size = Pt(14)
        run2.font.color.rgb = WHITE

    # Right Box (Alert illustration layout)
    right_card = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.0), Inches(2.2), Inches(4.5), Inches(4.0))
    right_card.fill.solid()
    right_card.fill.fore_color.rgb = RGBColor(30, 20, 20) # Reddish-dark card
    right_card.line.color.rgb = RGBColor(239, 68, 68)
    right_card.line.width = Pt(1.5)
    
    tf_r = right_card.text_frame
    tf_r.word_wrap = True
    tf_r.margin_left = Inches(0.3)
    tf_r.margin_right = Inches(0.3)
    tf_r.margin_top = Inches(0.3)
    
    p_r1 = tf_r.paragraphs[0]
    p_r1.text = "⚠️ ความเสี่ยงของระบบย่อย"
    p_r1.font.name = 'Sarabun'
    p_r1.font.size = Pt(18)
    p_r1.font.bold = True
    p_r1.font.color.rgb = RGBColor(239, 68, 68)
    p_r1.space_after = Pt(14)
    
    p_r2 = tf_r.add_paragraph()
    p_r2.text = "เนื่องจากระบบของทุกหน่วยงานเชื่อมโยงอยู่ในเน็ตเวิร์กเดียวกัน ช่องโหว่ของคณะเล็ก ๆ เพียงแห่งเดียว อาจกลายเป็นประตูให้แฮกเกอร์เจาะเข้าถึงเครือข่ายส่วนกลางทั้งหมดได้\n\nระบบนี้จึงมุ่งแก้ปัญหาด้วยการสร้างแหล่งข้อมูลแห่งความจริงหนึ่งเดียว (Single Source of Truth) เพื่อยกระดับความปลอดภัยร่วมกัน"
    p_r2.font.name = 'Sarabun'
    p_r2.font.size = Pt(13)
    p_r2.font.color.rgb = MUTED_GRAY

    # ==================== SLIDE 4: SOLUTION ====================
    slide4 = prs.slides.add_slide(slide_layout)
    add_slide_header(slide4, "แก้ปัญหาได้อย่างไร (How we solve it)", "03. Solution")
    
    # Left bullets
    left_box4 = slide4.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(6.5), Inches(4.5))
    tf_l4 = left_box4.text_frame
    tf_l4.word_wrap = True
    
    sol_bullets = [
        ("รวมข้อมูลจากหลายแหล่ง: ", "รวบรวมช่องโหว่ความมั่นคงปลอดภัยผ่านช่องทางสด (SecurityScorecard API) และการอัปโหลดไฟล์รายงานแบบแมนนวล (CSV Import)"),
        ("อัลกอริทึมแมปปิ้งอัตโนมัติ: ", "พัฒนาระบบสืบค้นโดเมนย่อย (Matched Domain) ชี้เป้าระบุหน่วยงานรับผิดชอบตามระบบบัญชีรายชื่อคณะ 21 หน่วยงานได้โดยอัตโนมัติ"),
        ("แสดงผลเชิงตำแหน่ง: ", "แผนที่วิทยาเขต Leaflet.js ปักหมุดที่ตั้งแต่ละคณะด้วยโลโก้คณะและรหัสสีตามระดับเกรด เพื่อให้ระบุจุดเสี่ยงบนพิกัดแผนที่ได้ทันที"),
        ("ระบบจัดการความรุนแรง: ", "จัดหมวดหมู่แยกความเสี่ยง (Critical, High, Medium, Info) ให้ทีมไอทีจัดลำดับการอุดช่องโหว่ได้ถูกต้อง")
    ]
    
    for i, (bold_txt, regular_txt) in enumerate(sol_bullets):
        p = tf_l4.add_paragraph() if i > 0 else tf_l4.paragraphs[0]
        p.space_after = Pt(16)
        
        run1 = p.add_run()
        run1.text = "•  " + bold_txt
        run1.font.name = 'Sarabun'
        run1.font.size = Pt(14)
        run1.font.bold = True
        run1.font.color.rgb = KKU_ORANGE
        
        run2 = p.add_run()
        run2.text = regular_txt
        run2.font.name = 'Sarabun'
        run2.font.size = Pt(14)
        run2.font.color.rgb = WHITE

    # Right Cards
    y_pos4 = Inches(2.0)
    right_cards = [
        ("คัดกรองระดับความปลอดภัยและเกรด", "ใช้ระบบวิเคราะห์เกรด (A, B, C, D, F) แปลงข้อมูลเทคนิคที่ยุ่งยากให้เป็นระดับคุณภาพ เพื่อสร้างความเข้าใจที่ตรงกันในระดับผู้บริหารและฝ่ายปฏิบัติการ"),
        ("ความปลอดภัยของข้อมูลสูงสุด", "เก็บรักษา API Key และค่าสำคัญในการเชื่อมโยงภายนอกด้วยการเข้ารหัสข้อมูลที่รัดกุมในระดับฐานข้อมูล (AES-256-GCM) เพื่อป้องกันข้อมูลผู้ใช้และคีย์รั่วไหล")
    ]
    
    for i, (title, desc) in enumerate(right_cards):
        card = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.8), y_pos4 + i * Inches(2.2), Inches(4.7), Inches(1.9))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = KKU_ORANGE if i == 0 else GREEN
        card.line.width = Pt(1.5)
        
        tf_c = card.text_frame
        tf_c.word_wrap = True
        tf_c.margin_left = Inches(0.2)
        tf_c.margin_right = Inches(0.2)
        tf_c.margin_top = Inches(0.2)
        
        p_t = tf_c.paragraphs[0]
        p_t.text = title
        p_t.font.name = 'Sarabun'
        p_t.font.size = Pt(15)
        p_t.font.bold = True
        p_t.font.color.rgb = WHITE
        p_t.space_after = Pt(6)
        
        p_d = tf_c.add_paragraph()
        p_d.text = desc
        p_d.font.name = 'Sarabun'
        p_d.font.size = Pt(12)
        p_d.font.color.rgb = MUTED_GRAY

    # ==================== SLIDE 5: METHODOLOGY ====================
    slide5 = prs.slides.add_slide(slide_layout)
    add_slide_header(slide5, "ทำอย่างไร (How did we do it?)", "04. Technology")
    
    # Left Column (System Architecture)
    left_box5 = slide5.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(6.0), Inches(4.5))
    tf_l5 = left_box5.text_frame
    tf_l5.word_wrap = True
    
    p_l5_1 = tf_l5.paragraphs[0]
    p_l5_1.text = "สถาปัตยกรรมระบบ (Full-stack Monorepo)"
    p_l5_1.font.name = 'Sarabun'
    p_l5_1.font.size = Pt(18)
    p_l5_1.font.bold = True
    p_l5_1.font.color.rgb = KKU_ORANGE
    p_l5_1.space_after = Pt(10)
    
    p_l5_2 = tf_l5.add_paragraph()
    p_l5_2.text = "ระบบได้รับการออกแบบเป็น Microservices ที่ทำงานร่วมกันอย่างสมบูรณ์แบบบน Docker-compose เพื่อความง่ายต่อการพัฒนาและติดตั้งใช้จริงในเซิร์ฟเวอร์กลาง"
    p_l5_2.font.name = 'Sarabun'
    p_l5_2.font.size = Pt(13)
    p_l5_2.font.color.rgb = MUTED_GRAY
    p_l5_2.space_after = Pt(15)
    
    arch_bullets = [
        ("Frontend: ", "Next.js ทำงานร่วมกับ Leaflet.js ในส่วนการแสดงผลและแผนที่โต้ตอบ"),
        ("Backend: ", "Node.js (Express Framework) พร้อมระบบจัดการตารางผ่าน Prisma ORM"),
        ("Database: ", "PostgreSQL จัดเก็บข้อมูลอย่างเป็นโครงสร้าง ทนทาน และมีประสิทธิภาพ"),
        ("Python Processor: ", "สคริปต์ Python วิเคราะห์ คัดกรอง และประมวลผลข้อมูลช่องโหว่ความเสี่ยง"),
        ("Nginx Proxy: ", "Reverse Proxy ช่วยจัดการเส้นทางและเชื่อมต่ออย่างปลอดภัย (HTTPS)")
    ]
    
    for bold_txt, regular_txt in arch_bullets:
        p = tf_l5.add_paragraph()
        p.space_after = Pt(6)
        
        run1 = p.add_run()
        run1.text = "• " + bold_txt
        run1.font.name = 'Sarabun'
        run1.font.size = Pt(13)
        run1.font.bold = True
        run1.font.color.rgb = KKU_ORANGE
        
        run2 = p.add_run()
        run2.text = regular_txt
        run2.font.name = 'Sarabun'
        run2.font.size = Pt(13)
        run2.font.color.rgb = WHITE

    # Right Column (Tech Badges)
    right_box5 = slide5.shapes.add_textbox(Inches(7.3), Inches(2.0), Inches(5.2), Inches(4.5))
    tf_r5 = right_box5.text_frame
    tf_r5.word_wrap = True
    
    p_r5_1 = tf_r5.paragraphs[0]
    p_r5_1.text = "เทคโนโลยีและเครื่องมือหลักที่เลือกใช้"
    p_r5_1.font.name = 'Sarabun'
    p_r5_1.font.size = Pt(18)
    p_r5_1.font.bold = True
    p_r5_1.font.color.rgb = KKU_ORANGE
    p_r5_1.space_after = Pt(15)
    
    # Let's create some visual rectangles for the tech stack
    techs = [
        ("Next.js (React)", "Frontend & Dashboard Interface"),
        ("Leaflet.js", "Interactive Campus Geolocation Map"),
        ("Node.js / Express", "RESTful API Backend Service"),
        ("Prisma & PostgreSQL", "Database & Type-safe Database Control"),
        ("Python Processing", "Data Analysis & Parsing Scripts"),
        ("Docker-compose", "Production & Development Environment")
    ]
    
    for i, (tech, desc) in enumerate(techs):
        row = tf_r5.add_paragraph()
        row.space_after = Pt(10)
        
        run_t = row.add_run()
        run_t.text = f"[{tech}] "
        run_t.font.name = 'Sarabun'
        run_t.font.size = Pt(13)
        run_t.font.bold = True
        run_t.font.color.rgb = GREEN
        
        run_d = row.add_run()
        run_d.text = desc
        run_d.font.name = 'Sarabun'
        run_d.font.size = Pt(13)
        run_d.font.color.rgb = WHITE

    # ==================== SLIDE 6: RESULTS ====================
    slide6 = prs.slides.add_slide(slide_layout)
    add_slide_header(slide6, "ผลลัพธ์เป็นอย่างไร (Results)", "05. Outcomes")
    
    # Left Column (Outcomes text)
    left_box6 = slide6.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(6.8), Inches(4.5))
    tf_l6 = left_box6.text_frame
    tf_l6.word_wrap = True
    
    outcomes = [
        ("ศูนย์รวมสถานะความปลอดภัย: ", "ระบบแดชบอร์ดช่วยให้ผู้ดูแลระบบกลาง มองเห็นคะแนนความปลอดภัย ประวัติความคืบหน้า และประวัติย้อนหลังของหน่วยงานย่อย 21 แห่งได้ทั้งหมดในที่เดียว"),
        ("ลดเวลาแก้ปัญหาแบบเจาะจง: ", "เปลี่ยนขั้นตอนค้นหาเจ้าของเครื่องเซิร์ฟเวอร์แบบเดิม ให้เป็นการกรองและแจ้งเตือนผ่าน API/CSV สังเคราะห์และส่งงานต่อให้เจ้าหน้าที่ไอทีคณะได้รวดเร็ว"),
        ("สร้างความตระหนักรู้การป้องกันภัย: ", "ระดับเกรด A-F ช่วยแปลงรายงานเชิงเทคนิคของช่องโหว่ (เช่น DNS record ผิดพลาด, SSL certificate หมดอายุ) ให้เป็นการวัดผลระดับองค์กรที่ผู้บริหารเข้าใจทันที")
    ]
    
    for i, (bold_txt, regular_txt) in enumerate(outcomes):
        p = tf_l6.add_paragraph() if i > 0 else tf_l6.paragraphs[0]
        p.space_after = Pt(20)
        
        run1 = p.add_run()
        run1.text = "•  " + bold_txt
        run1.font.name = 'Sarabun'
        run1.font.size = Pt(14)
        run1.font.bold = True
        run1.font.color.rgb = KKU_ORANGE
        
        run2 = p.add_run()
        run2.text = regular_txt
        run2.font.name = 'Sarabun'
        run2.font.size = Pt(14)
        run2.font.color.rgb = WHITE

    # Right Column (Grade Visual Card)
    grade_card = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.2), Inches(2.2), Inches(4.3), Inches(3.8))
    grade_card.fill.solid()
    grade_card.fill.fore_color.rgb = CARD_BG
    grade_card.line.color.rgb = GREEN
    grade_card.line.width = Pt(2)
    
    tf_gc = grade_card.text_frame
    tf_gc.word_wrap = True
    tf_gc.margin_left = Inches(0.3)
    tf_gc.margin_right = Inches(0.3)
    tf_gc.margin_top = Inches(0.4)
    
    p_g1 = tf_gc.paragraphs[0]
    p_g1.text = "A - F"
    p_g1.alignment = PP_ALIGN.CENTER
    p_g1.font.name = 'Sarabun'
    p_g1.font.size = Pt(48)
    p_g1.font.bold = True
    p_g1.font.color.rgb = GREEN
    p_g1.space_after = Pt(15)
    
    p_g2 = tf_gc.add_paragraph()
    p_g2.text = "ระบบการแปลคะแนนเชิงคุณภาพ"
    p_g2.alignment = PP_ALIGN.CENTER
    p_g2.font.name = 'Sarabun'
    p_g2.font.size = Pt(16)
    p_g2.font.bold = True
    p_g2.font.color.rgb = WHITE
    p_g2.space_after = Pt(10)
    
    p_g3 = tf_gc.add_paragraph()
    p_g3.text = "เปลี่ยนข้อมูลภัยคุกคามและช่องโหว่ความเสี่ยงทางเทคนิคที่เข้าใจยาก ให้กลายตัวเกรดสีความปลอดภัยที่แสดงผลบนหมุดแผนที่ เพื่อการสั่งการและแก้ไขอย่างเป็นระบบ"
    p_g3.alignment = PP_ALIGN.CENTER
    p_g3.font.name = 'Sarabun'
    p_g3.font.size = Pt(12)
    p_g3.font.color.rgb = MUTED_GRAY

    # ==================== SLIDE 7: PROGRESS ====================
    slide7 = prs.slides.add_slide(slide_layout)
    add_slide_header(slide7, "ตอนนี้งานไปถึงไหนแล้ว (Progress)", "06. Current Status")
    
    # Left Column: Progress bars
    left_box7 = slide7.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(6.0), Inches(4.5))
    tf_l7 = left_box7.text_frame
    tf_l7.word_wrap = True
    
    p_p1 = tf_l7.paragraphs[0]
    p_p1.text = "ความคืบหน้าการพัฒนาซอฟต์แวร์ย่อย"
    p_p1.font.name = 'Sarabun'
    p_p1.font.size = Pt(18)
    p_p1.font.bold = True
    p_p1.font.color.rgb = KKU_ORANGE
    p_p1.space_after = Pt(20)
    
    progresses = [
        ("Frontend Development (Next.js & Maps)", "85%", "ระบบ Dashboard, กราฟเชิงเปรียบเทียบ และ Interactive Map บนพิกัดดาวเทียมเสร็จสิ้นเป็นส่วนใหญ่"),
        ("Backend API (Express & Prisma)", "90%", "ระบบ API เชื่อมต่อฐานข้อมูล, ดึงข้อมูลจาก API ภายนอก และระบบเข้ารหัสรักษาความปลอดภัยสมบูรณ์"),
        ("Data Processing (Python Scripts)", "75%", "การเขียนสคริปต์จัดระเบียบข้อมูลและแปลข้อมูลช่องโหว่เชิงตำแหน่ง มีการสร้างสคริปต์หลักแล้ว")
    ]
    
    for title, pct, desc in progresses:
        p_t = tf_l7.add_paragraph()
        p_t.space_before = Pt(10)
        run_title = p_t.add_run()
        run_title.text = f"{title} : "
        run_title.font.name = 'Sarabun'
        run_title.font.size = Pt(13)
        run_title.font.bold = True
        run_title.font.color.rgb = WHITE
        
        run_pct = p_t.add_run()
        run_pct.text = pct
        run_pct.font.name = 'Sarabun'
        run_pct.font.size = Pt(13)
        run_pct.font.bold = True
        run_pct.font.color.rgb = KKU_ORANGE
        
        p_d = tf_l7.add_paragraph()
        p_d.space_after = Pt(10)
        run_desc = p_d.add_run()
        run_desc.text = desc
        run_desc.font.name = 'Sarabun'
        run_desc.font.size = Pt(11)
        run_desc.font.color.rgb = MUTED_GRAY

    # Right Column: Updates
    right_card7 = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.5), Inches(2.2), Inches(5.0), Inches(4.0))
    right_card7.fill.solid()
    right_card7.fill.fore_color.rgb = CARD_BG
    right_card7.line.color.rgb = KKU_ORANGE
    right_card7.line.width = Pt(1.5)
    
    tf_rc7 = right_card7.text_frame
    tf_rc7.word_wrap = True
    tf_rc7.margin_left = Inches(0.3)
    tf_rc7.margin_right = Inches(0.3)
    tf_rc7.margin_top = Inches(0.3)
    
    p_u1 = tf_rc7.paragraphs[0]
    p_u1.text = "การปรับปรุงล่าสุด (Recent Updates)"
    p_u1.font.name = 'Sarabun'
    p_u1.font.size = Pt(16)
    p_u1.font.bold = True
    p_u1.font.color.rgb = KKU_ORANGE
    p_u1.space_after = Pt(12)
    
    updates = [
        "ซ่อนแผนผังโครงข่ายหลักในหน้า Dashboard ชั่วคราวเพื่อความเรียบง่ายและเพิ่มพื้นที่แสดงผล",
        "ทำระบบข้ามขั้นตอนการล็อกอินแบบชั่วคราว (Login Bypass) เพื่อให้ผู้พัฒนาและผู้นำเสนอเข้าใช้งานได้ไวขึ้น",
        "รวมและจัดกลุ่มหน่วยงานย่อย (Organizations) เพื่อการแสดงสถานะโดยรวมที่กระชับขึ้น",
        "ปรับแต่งระบบ Fetcher และ Daily Cron Job เพื่อซิงค์ข้อมูลล่าสุดผ่าน API ของผู้ให้บริการประเมิน"
    ]
    
    for u in updates:
        p_u = tf_rc7.add_paragraph()
        p_u.space_after = Pt(8)
        
        run_bullet = p_u.add_run()
        run_bullet.text = "• "
        run_bullet.font.name = 'Sarabun'
        run_bullet.font.size = Pt(12)
        run_bullet.font.bold = True
        run_bullet.font.color.rgb = KKU_ORANGE
        
        run_text = p_u.add_run()
        run_text.text = u
        run_text.font.name = 'Sarabun'
        run_text.font.size = Pt(12)
        run_text.font.color.rgb = WHITE

    # Save presentation
    output_filename = "presentation.pptx"
    prs.save(output_filename)
    print(f"Presentation successfully saved as {output_filename}")

if __name__ == "__main__":
    create_presentation()
